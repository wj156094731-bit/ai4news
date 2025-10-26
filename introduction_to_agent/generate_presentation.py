#!/usr/bin/env python3
"""
PowerPoint Generator from Markdown
==================================

This script reads the AI_Agent_Capabilities_Demo.md file and generates
a professional PowerPoint presentation with the content.

Requirements:
- python-pptx library
- markdown parsing capabilities

Usage:
    python generate_presentation.py

Output:
    AI_Agent_Capabilities_Presentation.pptx
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class MarkdownToPowerPoint:
    def __init__(self, markdown_file, output_file):
        self.markdown_file = markdown_file
        self.output_file = output_file
        self.prs = Presentation()
        self.setup_slide_layout()
        
    def setup_slide_layout(self):
        """Set up the slide layout and theme"""
        # Use a clean, professional layout
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        return slide_layout
    
    def read_markdown_file(self):
        """Read and parse the markdown file"""
        with open(self.markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()
        return content
    
    def parse_sections(self, content):
        """Parse markdown content into sections"""
        sections = []
        
        # Split by main headers (##)
        parts = re.split(r'\n## ', content)
        
        for i, part in enumerate(parts):
            if i == 0:
                # Title section
                lines = part.strip().split('\n')
                title = lines[0].replace('# ', '')
                sections.append({
                    'type': 'title',
                    'title': title,
                    'content': '\n'.join(lines[1:]).strip()
                })
            else:
                # Regular sections
                lines = part.strip().split('\n')
                if lines:
                    title = lines[0]
                    content = '\n'.join(lines[1:]).strip()
                    sections.append({
                        'type': 'section',
                        'title': title,
                        'content': content
                    })
        
        return sections
    
    def create_title_slide(self, title, subtitle):
        """Create the title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
        
        return slide
    
    def create_content_slide(self, title, content):
        """Create a content slide with title and bullet points"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Parse content into bullet points
        lines = content.split('\n')
        current_bullet = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Handle different types of content
            if line.startswith('### '):
                # Subsection header
                p = text_frame.add_paragraph()
                p.text = line.replace('### ', '')
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)
                current_bullet = None
                
            elif line.startswith('**') and line.endswith('**'):
                # Bold text
                p = text_frame.add_paragraph()
                p.text = line.replace('**', '')
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(51, 51, 51)
                current_bullet = None
                
            elif line.startswith('- '):
                # Bullet point
                p = text_frame.add_paragraph()
                p.text = line.replace('- ', '')
                p.font.size = Pt(14)
                p.level = 0
                current_bullet = p
                
            elif line.startswith('  - '):
                # Sub-bullet point
                if current_bullet:
                    p = text_frame.add_paragraph()
                    p.text = line.replace('  - ', '')
                    p.font.size = Pt(12)
                    p.level = 1
                else:
                    p = text_frame.add_paragraph()
                    p.text = line.replace('  - ', '')
                    p.font.size = Pt(14)
                    p.level = 0
                    current_bullet = p
                    
            elif line.startswith('```'):
                # Code block - skip for now
                continue
                
            elif line.startswith('**What this means**'):
                # Special formatting for "What this means"
                p = text_frame.add_paragraph()
                p.text = line.replace('**', '')
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(204, 0, 0)
                current_bullet = None
                
            elif line and not line.startswith('---'):
                # Regular text
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                current_bullet = None
        
        return slide
    
    def create_example_slide(self, title, content):
        """Create a slide specifically for examples"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Content with special formatting for examples
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line or line.startswith('---'):
                continue
                
            if line.startswith('**Task**:') or line.startswith('**Process**:') or line.startswith('**Result**:'):
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)
                
            elif line.startswith('- '):
                p = text_frame.add_paragraph()
                p.text = line.replace('- ', '')
                p.font.size = Pt(14)
                p.level = 0
                
            elif line.startswith('  - '):
                p = text_frame.add_paragraph()
                p.text = line.replace('  - ', '')
                p.font.size = Pt(12)
                p.level = 1
                
            elif line.startswith('**What I can do**:'):
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(204, 0, 0)
        
        return slide
    
    def create_workflow_slide(self, title, content):
        """Create a slide for workflow diagrams"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line or line.startswith('---'):
                continue
                
            if line.startswith('```'):
                continue
            elif line.startswith('**Step '):
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)
                
            elif line.startswith('- '):
                p = text_frame.add_paragraph()
                p.text = line.replace('- ', '')
                p.font.size = Pt(14)
                p.level = 0
                
            elif line.startswith('**Result**:'):
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 153, 0)
        
        return slide
    
    def create_summary_slide(self, title, content):
        """Create a summary slide with key points"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line or line.startswith('---'):
                continue
                
            if line.startswith('### '):
                p = text_frame.add_paragraph()
                p.text = line.replace('### ', '')
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)
                
            elif line.startswith('- '):
                p = text_frame.add_paragraph()
                p.text = line.replace('- ', '')
                p.font.size = Pt(14)
                p.level = 0
                
            elif line.startswith('  - '):
                p = text_frame.add_paragraph()
                p.text = line.replace('  - ', '')
                p.font.size = Pt(12)
                p.level = 1
        
        return slide
    
    def generate_presentation(self):
        """Generate the complete PowerPoint presentation"""
        print("Reading markdown file...")
        content = self.read_markdown_file()
        
        print("Parsing content into sections...")
        sections = self.parse_sections(content)
        
        print("Creating slides...")
        
        # Create title slide
        title_section = sections[0]
        self.create_title_slide(
            title_section['title'],
            "HKU AI Agents Education Workshop - November 2025"
        )
        
        # Create overview slide
        overview_content = """
        This presentation demonstrates 5 unique capabilities that make AI agents 
        particularly valuable for research students:
        
        • Direct File System Operations & Code Execution
        • Multi-Tool Integration & Workflow Automation  
        • Real-Time Web Research & Information Synthesis
        • Interactive Problem-Solving & Iterative Development
        • Context-Aware Research Assistance
        
        Each capability includes practical examples from workshop materials.
        """
        
        self.create_content_slide("Overview", overview_content)
        
        # Process each main section
        for section in sections[1:]:
            title = section['title']
            content = section['content']
            
            # Clean up title (remove emojis and formatting)
            clean_title = re.sub(r'[^\w\s-]', '', title).strip()
            
            if 'Direct File System Operations' in title:
                self.create_content_slide(clean_title, content)
                
            elif 'Multi-Tool Integration' in title:
                self.create_workflow_slide(clean_title, content)
                
            elif 'Real-Time Web Research' in title:
                self.create_example_slide(clean_title, content)
                
            elif 'Interactive Problem-Solving' in title:
                self.create_content_slide(clean_title, content)
                
            elif 'Context-Aware Research' in title:
                self.create_example_slide(clean_title, content)
                
            elif 'Practical Workshop Examples' in title:
                self.create_workflow_slide(clean_title, content)
                
            elif 'Key Advantages' in title:
                self.create_summary_slide(clean_title, content)
                
            elif 'Workshop File Structure' in title:
                self.create_content_slide(clean_title, content)
                
            elif 'Next Steps' in title:
                self.create_summary_slide(clean_title, content)
        
        # Create conclusion slide
        conclusion_content = """
        Key Takeaway:
        
        AI agents can DO things, not just tell you how to do things.
        
        This makes them particularly valuable for hands-on research tasks 
        that require:
        
        • File manipulation and code execution
        • Tool integration across different software ecosystems
        • Real-time problem solving and debugging
        • Context-aware research assistance
        • Interactive workflow development
        
        Ready to explore AI agents in your research?
        """
        
        self.create_content_slide("Conclusion", conclusion_content)
        
        print(f"Saving presentation to {self.output_file}...")
        self.prs.save(self.output_file)
        print("Presentation generated successfully!")
        
        return self.output_file

def main():
    """Main function to generate the PowerPoint presentation"""
    
    # File paths
    markdown_file = "/Users/simonwang/Documents/Usage/HKUworkshop/agent4hku/introduction_to_agent/AI_Agent_Capabilities_Demo.md"
    output_file = "/Users/simonwang/Documents/Usage/HKUworkshop/agent4hku/introduction_to_agent/AI_Agent_Capabilities_Presentation.pptx"
    
    # Check if markdown file exists
    if not Path(markdown_file).exists():
        print(f"Error: Markdown file not found at {markdown_file}")
        return
    
    print("=" * 80)
    print("AI Agent Capabilities - PowerPoint Generator")
    print("=" * 80)
    print()
    
    try:
        # Create the presentation generator
        generator = MarkdownToPowerPoint(markdown_file, output_file)
        
        # Generate the presentation
        output_path = generator.generate_presentation()
        
        print()
        print("=" * 80)
        print("SUCCESS!")
        print("=" * 80)
        print(f"PowerPoint presentation created: {output_path}")
        print()
        print("The presentation includes:")
        print("• Title slide with overview")
        print("• 5 detailed capability sections")
        print("• Practical workshop examples")
        print("• Key advantages for research students")
        print("• Workshop file structure")
        print("• Next steps and conclusion")
        print()
        print("You can now open the PowerPoint file to view the presentation!")
        
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        print("\nMake sure you have the python-pptx library installed:")
        print("pip install python-pptx")

if __name__ == "__main__":
    main()

